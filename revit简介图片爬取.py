import os
import requests
from fake_useragent import UserAgent
from pptx import Presentation
from pptx.util import Inches


#import tempfile
# with tempfile.TemporaryDirectory() as tmpdirname:
#     print('创建临时目录', tmpdirname)


def mkdir(path):
    isExists = os.path.exists(path)
    if not isExists:
        os.makedirs(path)
        print(path+' 创建成功')
        return True
    else:
        print(path+' 目录已存在')
        return False


def get_content(url):
    ua = UserAgent()
    headers = {'User-Agent': ua.random}
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        return response.content
    else:
        return


prs = Presentation()
left, top, width, height = Inches(0), Inches(0), Inches(10), Inches(9)

mkpath = os.getcwd() + '\\Capture'
mkdir(mkpath)
urls = ['https://s3.ananas.chaoxing.com/doc/c3/d8/22/80995a297a2620fc245b16b1e57521b9/thumb/{}.png'.format(i+1) for i in range(55)]

for l in urls:
    img_path = mkpath + '\\' + l.split('/')[-1]
    with open(img_path, 'wb') as f:
        f.write(get_content(l))
    blank_slide_layout = prs.slide_layouts[3]                           
    slide = prs.slides.add_slide(blank_slide_layout)
    pic= slide.shapes.add_picture(img_path, left, top, width, height) 

  
prs.save('test.pptx')